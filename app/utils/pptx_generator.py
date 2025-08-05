from io import StringIO
from pptx import Presentation
from pathlib import Path
import os
import zipfile
from typing import Dict, Optional
import tempfile
from pptx.oxml import parse_xml
from pptx.util import Pt
from pptx.dml.color import RGBColor
from app.utils.drive_uploader import upload_file_to_drive, download_file_from_drive  
import re
import requests
import csv
from app.utils.generate_email_templates import generate_email_template

class PPTXGenerator:
    def __init__(self, variables: dict, form_id: int, files: Optional[dict] = None):
        self.variables = variables
        self.temp_dir = Path(tempfile.mkdtemp())
        self.form_id = form_id
        self.files = files

        plano = self.variables.get("plano")
        if plano:
            self.variables.update(self.get_features_by_plan(plano))

        planoold = self.variables.get("planoold")
        if planoold:
            features = self.get_features_by_plan(planoold)
            planoolddetalhes = {
                "descplanold": features.get("descplan"),
                "dimensionamentoold": features.get("dimensionamento")
            }
            self.variables.update(planoolddetalhes)

    # TODO Ajustar para ficar mais rápido. Quando eu tiver que buscar mais de uma informação, então eu devo retornar tudo em uma única consulta, para evitar gargalos.       
    @staticmethod
    def get_features_by_plan(plano: str) -> Dict[str, list]:
            sheet_id = os.getenv("SHEET_ID")
            sheet_name = os.getenv("SHEET_NAME", "Planos")
            url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/gviz/tq?tqx=out:csv&sheet={sheet_name}"

            response = requests.get(url)
            response.raise_for_status()

            # Ler o CSV
            descplan = []
            dimensionamento = []

            reader = csv.DictReader(StringIO(response.text))

            for row in reader:
                plano_csv = row.get("plano", "").strip().lower()
                if plano_csv != plano.lower():
                    continue

                # Preencher descplan e dimensionamento (ambos são opcionais)
                if row.get("descplan"):
                    desc = row["descplan"].strip()
                    if desc != "-":
                        descplan.append(desc)
                if row.get("dimensionamento"):
                    dimen = row["dimensionamento"].strip()
                    if dimen != "-":
                        dimensionamento.append(dimen)

            dic = {
                "descplan": descplan,
                "dimensionamento": dimensionamento
            }

            print(dic)
            return dic

    async def generate(self) -> Path:
        print(self.form_id)
        
        # Baixa o template do Drive
        template_drive_filename = f"{self.form_id}.pptx"
        template_path_str = await download_file_from_drive(template_drive_filename)
        template_path = Path(template_path_str)

        # Abre a apresentação
        prs = Presentation(template_path)
        
        # Manipula os dados
        if self.files:
            await self._handle_file_uploads(prs)
        self._replace_text_variables(prs)

        # Salva a apresentação final
        output_path = self.temp_dir / "Proposta e-Auditoria.pptx"
        prs.save(str(output_path))

        return output_path

    def _replace_text_variables(self, prs: Presentation):
        placeholders_slide = set()

        def substituir_placeholders(texto: str, dados: dict) -> str:
            encontrados = set(re.findall(r"\{\{(\w+)\}\}", texto))
            placeholders_slide.update(encontrados)

            for chave in encontrados:
                if chave not in dados:
                    print(f"⚠️  Placeholder '{{{{{chave}}}}}' está no slide mas NÃO foi encontrado em self.variables.")

            resultado = re.sub(r"\{\{(\w+)\}\}", lambda m: str(dados.get(m.group(1), "")), texto)

            return resultado

        def processar_text_frame(text_frame):
            if text_frame is None:
                return

            for paragrafo in text_frame.paragraphs:
                texto_original = paragrafo.text or ""
                estilo_base = None

                if paragrafo.runs:
                    run_base = paragrafo.runs[0]
                    try:
                        cor = run_base.font.color.rgb
                    except AttributeError:
                        cor = RGBColor(255, 255, 255)

                    estilo_base = {
                        "size": run_base.font.size,
                        "bold": run_base.font.bold,
                        "italic": run_base.font.italic,
                        "color": cor,
                        "name": run_base.font.name
                    }

                is_descplan = "{{descplan}}" in texto_original and isinstance(self.variables.get("descplan"), list)
                is_dimensionamento = "{{dimensionamento}}" in texto_original and isinstance(self.variables.get("dimensionamento"), list)
                is_descplanold = "{{descplanold}}" in texto_original and isinstance(self.variables.get("descplanold"), list)
                is_dimensionamentoold = "{{dimensionamentoold}}" in texto_original and isinstance(self.variables.get("dimensionamentoold"), list) 
                is_adicionais = "{{aditivos}}" in texto_original and "aditivos" in self.variables

                if is_descplan or is_dimensionamento or is_descplanold or is_dimensionamentoold or is_adicionais:
                    paragrafo.clear()
                    if is_descplan:
                        lista = self.variables["descplan"]
                    elif is_dimensionamento:
                        lista = self.variables["dimensionamento"]
                    elif is_descplanold:
                        lista = self.variables["descplanold"]
                    elif is_dimensionamentoold:
                        lista = self.variables["dimensionamentoold"]
                    elif is_adicionais:
                        lista = self.variables.get("aditivos", [])

                        if not lista:
                            texto_original = ''.join(run.text for run in paragrafo.runs)
                            paragrafo.clear()
                            run = paragrafo.add_run()
                            run.text = "Sem aditivos"

                            if estilo_base:
                                run.font.size = estilo_base["size"]
                                run.font.bold = estilo_base["bold"]
                                run.font.italic = estilo_base["italic"]
                                run.font.name = estilo_base["name"]
                                if estilo_base["color"]:
                                    run.font.color.rgb = estilo_base["color"]

                        else:
                            for index, item in enumerate(lista):
                                nome = item.get("nome", "").strip()
                                valor = str(item.get("valor", "")).strip()
                                texto = f"{nome}: {valor}"

                                if (index > 0):
                                    paragrafo = text_frame.add_paragraph()

                                paragrafo.clear()
                                run = paragrafo.add_run()
                                run.text = "◦" + u"\u00A0\u00A0" + texto
                                paragrafo.level = 0
                                if estilo_base:
                                    run.font.size = estilo_base["size"]
                                    run.font.bold = estilo_base["bold"]
                                    run.font.italic = estilo_base["italic"]
                                    run.font.name = estilo_base["name"]
                                    run.font.color.rgb = RGBColor(255, 255, 255)
                                paragrafo._element.set("lvl", "0")
                                paragrafo.bullet_indent = Pt(9)
                                paragrafo.margin_left = Pt(18)

                    for index, item in enumerate(lista):
                        texto = item
                
                        if is_adicionais:
                            continue
                        
                        if (index > 0):
                            paragrafo = text_frame.add_paragraph()

                        run = paragrafo.add_run()
                        run.text = "◦" + u"\u00A0\u00A0" + texto
                        paragrafo.level = 0
                        if estilo_base:
                            run.font.size = estilo_base["size"]
                            run.font.bold = estilo_base["bold"]
                            run.font.italic = estilo_base["italic"]
                            run.font.name = estilo_base["name"]
                            run.font.color.rgb = RGBColor(255, 255, 255)
                        paragrafo._element.set("lvl", "0")
                        paragrafo.bullet_indent = Pt(9)     # Bullet alinhado
                        paragrafo.margin_left = Pt(18)      # Texto mais afastado do bullet
                else:
                    texto_original = ''.join(run.text for run in paragrafo.runs)
                    novo_texto = substituir_placeholders(texto_original, self.variables)

                    if novo_texto != texto_original:
                        paragrafo.clear()
                        run = paragrafo.add_run()
                        
                        if "drive.google.com" in novo_texto:
                            run.text = "Baixar arquivo aqui"
                            run.hyperlink.address = novo_texto  # mantém o link clicável
                        else:
                            run.text = novo_texto

                        if estilo_base:
                            run.font.size = estilo_base["size"]
                            run.font.bold = estilo_base["bold"]
                            run.font.italic = estilo_base["italic"]
                            run.font.name = estilo_base["name"]
                            if estilo_base["color"]:
                                run.font.color.rgb = estilo_base["color"]

        def processar_tabela(tabela):
            for linha in tabela.rows:
                for celula in linha.cells:
                    celula.text = substituir_placeholders(celula.text, self.variables)
                    
        def processar_grupo(grupo):
            for shape in grupo.shapes:
                processar_shape(shape)

        def processar_shape(shape):
            if shape.shape_type == 6:
                processar_grupo(shape)
            elif shape.shape_type == 19:
                processar_tabela(shape.table)
            elif shape.has_text_frame:
                processar_text_frame(shape.text_frame)

        # Processa todos os slides
        for slide in prs.slides:
            for shape in slide.shapes:
                processar_shape(shape)

        # Validação final
        def validar_variaveis(placeholders_slide: set, variaveis: dict):
            faltando = placeholders_slide - set(variaveis.keys())
            if faltando:
                print("❌ ERRO: As seguintes variáveis estão ausentes em self.variables:")
                for chave in faltando:
                    print(f" - {{{{{chave}}}}}")

        validar_variaveis(placeholders_slide, self.variables)

        def processar_tabela(tabela):
            for linha in tabela.rows:
                for celula in linha.cells:
                    celula.text = substituir_placeholders(celula.text, self.variables)

        def processar_grupo(grupo):
            for shape in grupo.shapes:
                processar_shape(shape)

        def processar_shape(shape):
            if shape.shape_type == 6:
                processar_grupo(shape)
            elif shape.shape_type == 19:
                processar_tabela(shape.table)
            elif shape.has_text_frame:
                processar_text_frame(shape.text_frame)

        for slide in prs.slides:
            for shape in slide.shapes:
                processar_shape(shape)

    async def _handle_file_uploads(self, prs):
        for key, file in self.files.items():
            if file is None:
                continue

            file_path = self.temp_dir / file.filename
            with open(file_path, "wb") as f:
                content = await file.read()
                f.write(content)

            link = await upload_file_to_drive(str(file_path), file.filename)

            # Substitui o valor da variável (ex: down1, down2...) pelo link do Drive
            self.variables[key] = link

    def generate_email_template(self) -> Path:
        email_template = generate_email_template(self.form_id, self.variables)
        email_path = self.temp_dir / "email_template.txt"
        email_path.write_text(email_template)
        
        return email_path

    def create_zip(self, pptx_path: Path, pdf_path: Path, email_path: Path) -> Path:
        zip_path = self.temp_dir / "proposta e-Auditoria.zip"
        with zipfile.ZipFile(zip_path, 'w') as zip_file:
            zip_file.write(pptx_path, "proposta e-Auditoria.pptx")
            zip_file.write(pdf_path, "proposta e-Auditoria.pdf")
            zip_file.write(email_path, "email proposta e-Auditoria.txt")
            if self.files:
                for file in self.files.values():
                    if file is None:
                        continue
                    file_path = self.temp_dir / file.filename
                    if file_path.exists():
                        zip_file.write(file_path, f"anexos/{file.filename}")
        return zip_path